# Роля на модула: PPTX renderer. Получава вече изчислен semantic layout и го превръща в реални PowerPoint shapes.
# Чети коментарите като обяснение на причината за кода и връзката му със следващия слой, а не като буквален превод на Python синтаксиса.
from __future__ import annotations

from pathlib import Path
from typing import Any

from PIL import Image
from pptx import Presentation as PptxPresentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

from app.semantic.contracts import (
    Alignment,
    LayoutedPresentationDocument,
    LayoutedSlide,
    LayoutElement,
    LayoutElementKind,
    ThemeDefinition,
)

# `SLIDE_WIDTH` пази резултата от `Inches`, за да бъде проверен или използван в следващите стъпки вместо операцията да се повтори.
SLIDE_WIDTH = Inches(13.333)
# `SLIDE_HEIGHT` пази резултата от `Inches`, за да бъде проверен или използван в следващите стъпки вместо операцията да се повтори.
SLIDE_HEIGHT = Inches(7.5)
BLANK_LAYOUT_INDEX = 6


def _rgb(value: str) -> RGBColor:
    # Роля в pipeline-а: Това е вътрешна помощна стъпка: обработва стъпката `rgb` като отделна отговорност, така че caller-ът да използва резултата без да познава вътрешните проверки и междинни стойности.
    # Входът идва през `value` (str); имената показват каква част от контекста е собственост на тази стъпка.
    # Основните преходи навън са към `value.strip().lstrip`, `RGBColor`; така се вижда кои отговорности функцията делегира.
    # Типовете в сигнатурата документират договора за caller-а и позволяват грешки да се хващат преди runtime.
    # Изходен договор: `RGBColor`. Резултатът е част от последния rendering/export етап и вече е близо до крайния PPTX/PDF файл.
    # `cleaned` пази резултата от `value.strip().lstrip`, за да бъде проверен или използван в следващите стъпки вместо операцията да се повтори.
    cleaned = value.strip().lstrip("#")
    # Това условие е decision point: `len(cleaned) == 3`.
    # При вярно условие се активира `''.join`; така този branch избира конкретна стратегия, а не просто проверява стойност.
    if len(cleaned) == 3:
        # `cleaned` пази резултата от `''.join`, за да бъде проверен или използван в следващите стъпки вместо операцията да се повтори.
        cleaned = "".join(char * 2 for char in cleaned)
    # Това условие е decision point: `len(cleaned) != 6 or not all((char in '0123456789abcdefABCDEF' for char in cleaned))`.
    # Това е guard clause: при вярно условие вече имаме достатъчно надежден резултат (`RGBColor(37, 99, 235)`) и прескачаме ненужната останала работа.
    if len(cleaned) != 6 or not all(char in "0123456789abcdefABCDEF" for char in cleaned):
        return RGBColor(37, 99, 235)
    return RGBColor(*(int(cleaned[index : index + 2], 16) for index in (0, 2, 4)))


PX_TO_INCH = 1.0 / 96.0


def _px(value: int) -> float:
    # Роля в pipeline-а: Това е вътрешна помощна стъпка: обработва стъпката `px` като отделна отговорност, така че caller-ът да използва резултата без да познава вътрешните проверки и междинни стойности.
    # Входът идва през `value` (int); имената показват каква част от контекста е собственост на тази стъпка.
    # Функцията работи основно с локални стойности и не делегира към други services.
    # Типовете в сигнатурата документират договора за caller-а и позволяват грешки да се хващат преди runtime.
    # Изходен договор: `float`. Резултатът е част от последния rendering/export етап и вече е близо до крайния PPTX/PDF файл.
    return value * PX_TO_INCH


def _component_style(theme: ThemeDefinition, component: str) -> dict:
    # Роля в pipeline-а: Това е вътрешна помощна стъпка: обработва стъпката `component_style` като отделна отговорност, така че caller-ът да използва резултата без да познава вътрешните проверки и междинни стойности.
    # Входът идва през `theme` (ThemeDefinition), `component` (str); имената показват каква част от контекста е собственост на тази стъпка.
    # Функцията работи основно с локални стойности и не делегира към други services.
    # Типовете в сигнатурата документират договора за caller-а и позволяват грешки да се хващат преди runtime.
    # Изходен договор: `dict`. Резултатът е част от последния rendering/export етап и вече е близо до крайния PPTX/PDF файл.
    return theme.tokens.component_styles.get(component, {})


def _add_fitted_picture(slide: Any, image_path: Path, left: int, top: int, width: int, height: int, fit: str) -> None:
    # Роля в pipeline-а: Това е вътрешна помощна стъпка: обработва стъпката `add_fitted_picture` като отделна отговорност, така че caller-ът да използва резултата без да познава вътрешните проверки и междинни стойности.
    # Входът идва през `slide` (Any), `image_path` (Path), `left` (int), `top` (int), `width` (int), `height` (int) и още параметри; имената показват каква част от контекста е собственост на тази стъпка.
    # Основните преходи навън са към `slide.shapes.add_picture`, `Image.open`, `Inches`, `_px`; така се вижда кои отговорности функцията делегира.
    # Типовете в сигнатурата документират договора за caller-а и позволяват грешки да се хващат преди runtime.
    # Изходен договор: функцията не връща нов обект; ефектът ѝ е промяна на подадено състояние, файл или външна услуга.
    with Image.open(image_path) as image:
        source_width, source_height = image.size
    source_ratio = source_width / max(source_height, 1)
    frame_ratio = width / max(height, 1)
    # Това условие е decision point: `fit == 'contain'`.
    # Това е guard clause: при вярно условие вече имаме достатъчно надежден резултат (`None`) и прескачаме ненужната останала работа.
    if fit == "contain":
        # Това условие е decision point: `source_ratio > frame_ratio`.
        # При вярно условие се променя текущото състояние, което влияе на следващите стъпки.
        if source_ratio > frame_ratio:
            render_width = width
            render_height = width / source_ratio
        else:
            render_height = height
            render_width = height * source_ratio
        slide.shapes.add_picture(
            str(image_path),
            Inches(_px(left + ((width - render_width) / 2))),
            Inches(_px(top + ((height - render_height) / 2))),
            Inches(_px(render_width)),
            Inches(_px(render_height)),
        )
        return
    # `picture` пази резултата от `slide.shapes.add_picture`, за да бъде проверен или използван в следващите стъпки вместо операцията да се повтори.
    picture = slide.shapes.add_picture(
        str(image_path), Inches(_px(left)), Inches(_px(top)), Inches(_px(width)), Inches(_px(height))
    )
    # Това условие е decision point: `source_ratio > frame_ratio`.
    # При вярно условие се активира `max`; така този branch избира конкретна стратегия, а не просто проверява стойност.
    if source_ratio > frame_ratio:
        crop = max(0.0, (1 - (frame_ratio / source_ratio)) / 2)
        picture.crop_left = crop
        picture.crop_right = crop
    # Това условие е decision point: `source_ratio < frame_ratio`.
    # При вярно условие се активира `max`; така този branch избира конкретна стратегия, а не просто проверява стойност.
    elif source_ratio < frame_ratio:
        crop = max(0.0, (1 - (source_ratio / frame_ratio)) / 2)
        picture.crop_top = crop
        picture.crop_bottom = crop


def _add_semantic_icon(slide: Any, icon: str, left: float, top: float, size: float, theme: ThemeDefinition) -> None:
    # Роля в pipeline-а: Това е вътрешна помощна стъпка: обработва стъпката `add_semantic_icon` като отделна отговорност, така че caller-ът да използва резултата без да познава вътрешните проверки и междинни стойности.
    # Входът идва през `slide` (Any), `icon` (str), `left` (float), `top` (float), `size` (float), `theme` (ThemeDefinition); имената показват каква част от контекста е собственост на тази стъпка.
    # Основните преходи навън са към `_rgb`, `slide.shapes.add_shape`, `container.fill.solid`, `Pt`; така се вижда кои отговорности функцията делегира.
    # Типовете в сигнатурата документират договора за caller-а и позволяват грешки да се хващат преди runtime.
    # Изходен договор: функцията не връща нов обект; ефектът ѝ е промяна на подадено състояние, файл или външна услуга.
    # `accent` пази резултата от `_rgb`, за да бъде проверен или използван в следващите стъпки вместо операцията да се повтори.
    accent = _rgb(theme.tokens.accent_primary)
    # `soft` пази резултата от `_rgb`, за да бъде проверен или използван в следващите стъпки вместо операцията да се повтори.
    soft = _rgb(theme.tokens.accent_secondary)

    def add_shape(kind: Any, x: float, y: float, width: float, height: float, *, filled: bool = False) -> Any:
        # `shape` пази резултата от `slide.shapes.add_shape`, за да бъде проверен или използван в следващите стъпки вместо операцията да се повтори.
        # Роля в pipeline-а: обработва стъпката `add_shape` като отделна отговорност, така че caller-ът да използва резултата без да познава вътрешните проверки и междинни стойности.
        # Входът идва през `kind` (Any), `x` (float), `y` (float), `width` (float), `height` (float), `filled` (bool); имената показват каква част от контекста е собственост на тази стъпка.
        # Основните преходи навън са към `slide.shapes.add_shape`, `Pt`, `Inches`, `shape.fill.solid`; така се вижда кои отговорности функцията делегира.
        # Типовете в сигнатурата документират договора за caller-а и позволяват грешки да се хващат преди runtime.
        # Изходен договор: `Any`. Резултатът е част от последния rendering/export етап и вече е близо до крайния PPTX/PDF файл.
        shape = slide.shapes.add_shape(kind, Inches(_px(x)), Inches(_px(y)), Inches(_px(width)), Inches(_px(height)))
        # Това условие е decision point: `filled`.
        # При вярно условие се активира `shape.fill.solid`; така този branch избира конкретна стратегия, а не просто проверява стойност.
        if filled:
            shape.fill.solid()
            shape.fill.fore_color.rgb = accent
        else:
            shape.fill.background()
        shape.line.color.rgb = accent
        # `shape.line.width` пази резултата от `Pt`, за да бъде проверен или използван в следващите стъпки вместо операцията да се повтори.
        shape.line.width = Pt(1)
        return shape

    def add_line(x1: float, y1: float, x2: float, y2: float) -> None:
        # `line` пази резултата от `slide.shapes.add_connector`, за да бъде проверен или използван в следващите стъпки вместо операцията да се повтори.
        # Роля в pipeline-а: обработва стъпката `add_line` като отделна отговорност, така че caller-ът да използва резултата без да познава вътрешните проверки и междинни стойности.
        # Входът идва през `x1` (float), `y1` (float), `x2` (float), `y2` (float); имената показват каква част от контекста е собственост на тази стъпка.
        # Основните преходи навън са към `slide.shapes.add_connector`, `Pt`, `Inches`, `_px`; така се вижда кои отговорности функцията делегира.
        # Типовете в сигнатурата документират договора за caller-а и позволяват грешки да се хващат преди runtime.
        # Изходен договор: функцията не връща нов обект; ефектът ѝ е промяна на подадено състояние, файл или външна услуга.
        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Inches(_px(x1)), Inches(_px(y1)), Inches(_px(x2)), Inches(_px(y2))
        )
        line.line.color.rgb = accent
        # `line.line.width` пази резултата от `Pt`, за да бъде проверен или използван в следващите стъпки вместо операцията да се повтори.
        line.line.width = Pt(1)

    # `container` пази резултата от `slide.shapes.add_shape`, за да бъде проверен или използван в следващите стъпки вместо операцията да се повтори.
    container = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(_px(left)), Inches(_px(top)), Inches(_px(size)), Inches(_px(size))
    )
    container.fill.solid()
    container.fill.fore_color.rgb = soft
    container.line.color.rgb = accent
    # `container.line.width` пази резултата от `Pt`, за да бъде проверен или използван в следващите стъпки вместо операцията да се повтори.
    container.line.width = Pt(0.8)
    x = left + (size * 0.23)
    y = top + (size * 0.23)
    inner = size * 0.54
    # Това условие е decision point: `icon == 'chart'`.
    # При вярно условие се активира `enumerate`; така този branch избира конкретна стратегия, а не просто проверява стойност.
    if icon == "chart":
        # Обхождаме `enumerate((0.42, 0.72, 1.0))` като `(index, factor)`, защото всеки елемент трябва да мине през една и съща pipeline стъпка.
        # Цикълът държи обработката еднаква за всеки елемент.
        for index, factor in enumerate((0.42, 0.72, 1.0)):
            bar_height = inner * factor
            add_shape(
                MSO_SHAPE.RECTANGLE,
                x + (index * inner * 0.34),
                y + inner - bar_height,
                inner * 0.16,
                bar_height,
                filled=True,
            )
    # Това условие е decision point: `icon == 'bolt'`.
    # При вярно условие се активира `add_shape`; така този branch избира конкретна стратегия, а не просто проверява стойност.
    elif icon == "bolt":
        add_shape(MSO_SHAPE.LIGHTNING_BOLT, x, y, inner, inner, filled=True)
    # Това условие е decision point: `icon == 'star'`.
    # При вярно условие се активира `add_shape`; така този branch избира конкретна стратегия, а не просто проверява стойност.
    elif icon == "star":
        add_shape(MSO_SHAPE.STAR_5_POINT, x, y, inner, inner, filled=True)
    # Това условие е decision point: `icon == 'heart'`.
    # При вярно условие се активира `add_shape`; така този branch избира конкретна стратегия, а не просто проверява стойност.
    elif icon == "heart":
        add_shape(MSO_SHAPE.HEART, x, y, inner, inner, filled=True)
    # Това условие е decision point: `icon in {'person', 'users'}`.
    # При вярно условие се активира `add_shape`; така този branch избира конкретна стратегия, а не просто проверява стойност.
    elif icon in {"person", "users"}:
        add_shape(MSO_SHAPE.OVAL, x + inner * 0.32, y, inner * 0.36, inner * 0.36, filled=True)
        add_shape(MSO_SHAPE.OVAL, x + inner * 0.14, y + inner * 0.42, inner * 0.72, inner * 0.5)
    # Това условие е decision point: `icon == 'clock'`.
    # При вярно условие се активира `add_shape`; така този branch избира конкретна стратегия, а не просто проверява стойност.
    elif icon == "clock":
        add_shape(MSO_SHAPE.OVAL, x, y, inner, inner)
        add_line(x + inner / 2, y + inner / 2, x + inner / 2, y + inner * 0.2)
        add_line(x + inner / 2, y + inner / 2, x + inner * 0.74, y + inner * 0.58)
    elif icon == "eye":
        add_shape(MSO_SHAPE.OVAL, x, y + inner * 0.16, inner, inner * 0.68)
        add_shape(MSO_SHAPE.OVAL, x + inner * 0.34, y + inner * 0.34, inner * 0.32, inner * 0.32, filled=True)
    elif icon == "book":
        add_shape(MSO_SHAPE.RECTANGLE, x, y + inner * 0.08, inner * 0.46, inner * 0.84)
        add_shape(MSO_SHAPE.RECTANGLE, x + inner * 0.54, y + inner * 0.08, inner * 0.46, inner * 0.84)
    elif icon in {"map", "route"}:
        add_line(x + inner * 0.12, y + inner * 0.82, x + inner * 0.48, y + inner * 0.28)
        add_line(x + inner * 0.48, y + inner * 0.28, x + inner * 0.88, y + inner * 0.68)
        # Обхождаме `((0.12, 0.82), (0.48, 0.28), (0.88, 0.68))` като `(px, py)`, защото всеки елемент трябва да мине през една и съща pipeline стъпка.
        # Цикълът държи обработката еднаква за всеки елемент.
        for px, py in ((0.12, 0.82), (0.48, 0.28), (0.88, 0.68)):
            add_shape(MSO_SHAPE.OVAL, x + inner * px - 2, y + inner * py - 2, 4, 4, filled=True)
    elif icon == "search":
        add_shape(MSO_SHAPE.OVAL, x, y, inner * 0.68, inner * 0.68)
        add_line(x + inner * 0.56, y + inner * 0.56, x + inner, y + inner)
    elif icon == "film":
        add_shape(MSO_SHAPE.RECTANGLE, x, y + inner * 0.12, inner, inner * 0.76)
        add_line(x + inner * 0.25, y + inner * 0.12, x + inner * 0.25, y + inner * 0.88)
        add_line(x + inner * 0.75, y + inner * 0.12, x + inner * 0.75, y + inner * 0.88)
    elif icon == "shield":
        add_shape(MSO_SHAPE.PENTAGON, x, y, inner, inner, filled=True)
    elif icon == "home":
        add_shape(MSO_SHAPE.PENTAGON, x, y, inner, inner)
        add_shape(MSO_SHAPE.RECTANGLE, x + inner * 0.38, y + inner * 0.55, inner * 0.24, inner * 0.45)
    elif icon == "idea":
        add_shape(MSO_SHAPE.SUN, x, y, inner, inner)
    else:
        add_shape(MSO_SHAPE.OVAL, x, y, inner, inner)
        add_shape(MSO_SHAPE.OVAL, x + inner * 0.24, y + inner * 0.24, inner * 0.52, inner * 0.52)


def _semantic_font_name(theme: ThemeDefinition, region: str, kind: LayoutElementKind) -> str:
    # Роля в pipeline-а: Това е вътрешна помощна стъпка: обработва стъпката `semantic_font_name` като отделна отговорност, така че caller-ът да използва резултата без да познава вътрешните проверки и междинни стойности.
    # Входът идва през `theme` (ThemeDefinition), `region` (str), `kind` (LayoutElementKind); имената показват каква част от контекста е собственост на тази стъпка.
    # Функцията работи основно с локални стойности и не делегира към други services.
    # Типовете в сигнатурата документират договора за caller-а и позволяват грешки да се хващат преди runtime.
    # Изходен договор: `str`. Резултатът е част от последния rendering/export етап и вече е близо до крайния PPTX/PDF файл.
    # Това условие е decision point: `kind in {LayoutElementKind.QUOTE, LayoutElementKind.STATISTIC}`.
    # Това е guard clause: при вярно условие вече имаме достатъчно надежден резултат (`theme.tokens.fonts.heading or theme.tokens.fonts.body`) и прескачаме ненужната останала работа.
    if kind in {LayoutElementKind.QUOTE, LayoutElementKind.STATISTIC}:
        return theme.tokens.fonts.heading or theme.tokens.fonts.body
    # Това условие е decision point: `any((token in region for token in ('title', 'quote', 'attribution', 'heading')))`.
    # Това е guard clause: при вярно условие вече имаме достатъчно надежден резултат (`theme.tokens.fonts.heading or theme.tokens.fonts.body`) и прескачаме ненужната останала работа.
    if any(token in region for token in ("title", "quote", "attribution", "heading")):
        return theme.tokens.fonts.heading or theme.tokens.fonts.body
    return theme.tokens.fonts.body or theme.tokens.fonts.heading


def _render_debug_label(slide: Any, element: LayoutElement, left: float, top: float, theme: ThemeDefinition) -> None:
    # Роля в pipeline-а: Това е вътрешна помощна стъпка: превежда semantic/layout информация към конкретни визуални обекти.
    # Входът идва през `slide` (Any), `element` (LayoutElement), `left` (float), `top` (float), `theme` (ThemeDefinition); имената показват каква част от контекста е собственост на тази стъпка.
    # Основните преходи навън са към `slide.shapes.add_textbox`, `frame.clear`, `frame.paragraphs[0].add_run`, `Pt`; така се вижда кои отговорности функцията делегира.
    # Типовете в сигнатурата документират договора за caller-а и позволяват грешки да се хващат преди runtime.
    # Изходен договор: функцията не връща нов обект; ефектът ѝ е промяна на подадено състояние, файл или външна услуга.
    # `box` пази резултата от `slide.shapes.add_textbox`, за да бъде проверен или използван в следващите стъпки вместо операцията да се повтори.
    box = slide.shapes.add_textbox(
        Inches(_px(left + 4)), Inches(_px(top + 4)), Inches(max(_px(min(element.width, 200)), 0.4)), Inches(0.22)
    )
    frame = box.text_frame
    frame.clear()
    # `run` пази резултата от `frame.paragraphs[0].add_run`, за да бъде проверен или използван в следващите стъпки вместо операцията да се повтори.
    run = frame.paragraphs[0].add_run()
    # `run.text` е нормализирано работно копие на текста; оригиналът остава непокътнат, а проверките стават върху предвидим формат.
    run.text = f"{element.region} · {element.x},{element.y} {element.width}x{element.height}"
    run.font.name = theme.tokens.fonts.mono or theme.tokens.fonts.body
    # `run.font.size` пази резултата от `Pt`, за да бъде проверен или използван в следващите стъпки вместо операцията да се повтори.
    run.font.size = Pt(7)
    # `run.font.color.rgb` пази резултата от `_rgb`, за да бъде проверен или използван в следващите стъпки вместо операцията да се повтори.
    run.font.color.rgb = _rgb(theme.tokens.accent_primary)


def _render_layout_element(
    slide: Any,
    element: LayoutElement,
    theme: ThemeDefinition,
    *,
    offset_x: int = 0,
    offset_y: int = 0,
    debug_mode: bool = False,
) -> None:
    # Роля в pipeline-а: Това е вътрешна помощна стъпка: превежда semantic/layout информация към конкретни визуални обекти.
    # Входът идва през `slide` (Any), `element` (LayoutElement), `theme` (ThemeDefinition), `offset_x` (int), `offset_y` (int), `debug_mode` (bool); имената показват каква част от контекста е собственост на тази стъпка.
    # Основните преходи навън са към `slide.shapes.add_shape`, `panel.fill.solid`, `_rgb`, `Pt`; така се вижда кои отговорности функцията делегира.
    # Типовете в сигнатурата документират договора за caller-а и позволяват грешки да се хващат преди runtime.
    # Изходен договор: функцията не връща нов обект; ефектът ѝ е промяна на подадено състояние, файл или външна услуга.
    left = offset_x + element.x
    top = offset_y + element.y
    width = element.width
    height = element.height

    # Това условие е decision point: `element.kind == LayoutElementKind.PANEL`.
    # При вярно условие се активира `slide.shapes.add_shape`; така този branch избира конкретна стратегия, а не просто проверява стойност.
    if element.kind == LayoutElementKind.PANEL:
        # `panel_shape` пази резултата от `_component_style(theme, 'panel').get`, за да бъде проверен или използван в следващите стъпки вместо операцията да се повтори.
        panel_shape = (
            MSO_SHAPE.RECTANGLE
            if _component_style(theme, "panel").get("style") == "square"
            else MSO_SHAPE.ROUNDED_RECTANGLE
        )
        # `panel` пази резултата от `slide.shapes.add_shape`, за да бъде проверен или използван в следващите стъпки вместо операцията да се повтори.
        panel = slide.shapes.add_shape(
            panel_shape, Inches(_px(left)), Inches(_px(top)), Inches(_px(width)), Inches(_px(height))
        )
        panel.fill.solid()
        # `panel.fill.fore_color.rgb` пази резултата от `_rgb`, за да бъде проверен или използван в следващите стъпки вместо операцията да се повтори.
        panel.fill.fore_color.rgb = _rgb(theme.tokens.surface)
        # `panel.line.color.rgb` пази резултата от `_rgb`, за да бъде проверен или използван в следващите стъпки вместо операцията да се повтори.
        panel.line.color.rgb = _rgb(theme.tokens.border)
        # `panel.line.width` пази резултата от `Pt`, за да бъде проверен или използван в следващите стъпки вместо операцията да се повтори.
        panel.line.width = Pt(1)
    # Това условие е decision point: `element.kind == LayoutElementKind.IMAGE`.
    # При вярно условие се активира `_component_style`; така този branch избира конкретна стратегия, а не просто проверява стойност.
    elif element.kind == LayoutElementKind.IMAGE:
        # `local_path` материализира резултата като локална файлова референция, която renderer-ът може да използва без нова мрежова заявка.
        local_path = element.content.get("local_path") if isinstance(element.content, dict) else None
        # `image_path` е локалният asset път, който renderer/exporter може да отвори без мрежова зависимост.
        image_path = Path(local_path) if local_path else None
        image_style = _component_style(theme, "image")
        inset = max(0, int(image_style.get("frame_inset") or 0))
        frame_shape = (
            MSO_SHAPE.RECTANGLE
            if _component_style(theme, "panel").get("style") == "square"
            else MSO_SHAPE.ROUNDED_RECTANGLE
        )
        frame_shape = slide.shapes.add_shape(
            frame_shape, Inches(_px(left)), Inches(_px(top)), Inches(_px(width)), Inches(_px(height))
        )
        frame_shape.fill.solid()
        frame_shape.fill.fore_color.rgb = _rgb(theme.tokens.background_alt)
        frame_shape.line.color.rgb = _rgb(theme.tokens.border)
        frame_shape.line.width = Pt(1)
        # Това условие е decision point: `image_path and image_path.exists()`.
        # При вярно условие се активира `_add_fitted_picture`; така този branch избира конкретна стратегия, а не просто проверява стойност.
        if image_path and image_path.exists():
            _add_fitted_picture(
                slide,
                image_path,
                left + inset,
                top + inset,
                max(1, width - (inset * 2)),
                max(1, height - (inset * 2)),
                str(element.content.get("fit") or image_style.get("fit") or "cover"),
            )
        else:
            # Keep unresolved image frames visually empty; prompt text is only
            # debugging metadata and should not appear as slide content.
            pass
    # Това условие е decision point: `element.kind != LayoutElementKind.BULLET_LIST`.
    # При вярно условие се активира `slide.shapes.add_textbox`; така този branch избира конкретна стратегия, а не просто проверява стойност.
    elif element.kind != LayoutElementKind.BULLET_LIST:
        text_left = left
        text_width = width
        text_top = top
        text_height = height
        # Това условие е decision point: `element.kind == LayoutElementKind.BULLET_ITEM`.
        # При вярно условие се активира `min`; така този branch избира конкретна стратегия, а не просто проверява стойност.
        if element.kind == LayoutElementKind.BULLET_ITEM:
            if _component_style(theme, "bullet").get("style") == "lines":
                line = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, Inches(_px(left)), Inches(_px(top)), Inches(_px(3)), Inches(_px(height))
                )
                line.fill.solid()
                line.fill.fore_color.rgb = _rgb(theme.tokens.accent_primary)
                line.line.fill.background()
            else:
                card = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(_px(left)),
                    Inches(_px(top)),
                    Inches(_px(width)),
                    Inches(_px(height)),
                )
                card.fill.solid()
                card.fill.fore_color.rgb = _rgb(theme.tokens.surface)
                card.line.color.rgb = _rgb(theme.tokens.border)
                card.line.width = Pt(0.8)
            marker_size = min(28, max(20, height - 12))
            _add_semantic_icon(
                slide,
                str(element.content.get("icon") or "target"),
                left + 10,
                top + ((height - marker_size) / 2),
                marker_size,
                theme,
            )
            text_left += 50
            text_width -= 58
            text_top += 5
            text_height -= 10
        textbox = slide.shapes.add_textbox(
            Inches(_px(text_left)),
            Inches(_px(text_top)),
            Inches(_px(max(1, text_width))),
            Inches(_px(max(1, text_height))),
        )
        frame = textbox.text_frame
        frame.word_wrap = element.wrap
        frame.auto_size = MSO_AUTO_SIZE.NONE
        frame.vertical_anchor = MSO_ANCHOR.TOP
        frame.clear()
        frame.margin_left = Inches(0.02)
        frame.margin_right = Inches(0.02)
        frame.margin_top = Inches(0.02)
        frame.margin_bottom = Inches(0.02)
        paragraph = frame.paragraphs[0]
        paragraph.alignment = {
            Alignment.START: PP_ALIGN.LEFT,
            Alignment.CENTER: PP_ALIGN.CENTER,
            Alignment.END: PP_ALIGN.RIGHT,
        }[element.align]
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(0)
        paragraph.line_spacing = min(element.line_height or 1.18, 1.18)
        run = paragraph.add_run()
        run.text = element.text or ""
        run.font.name = _semantic_font_name(theme, element.region, element.kind)
        # Това условие е decision point: `element.font_size`.
        # При вярно условие се активира `Pt`; така този branch избира конкретна стратегия, а не просто проверява стойност.
        if element.font_size:
            run.font.size = Pt(round(element.font_size * 0.75, 1))
        run.font.color.rgb = _rgb(
            theme.tokens.text_secondary
            if any(token in element.region for token in ("subtitle", "notes", "attribution"))
            else theme.tokens.text_primary
        )

    # Това условие е decision point: `debug_mode`.
    # При вярно условие се активира `slide.shapes.add_shape`; така този branch избира конкретна стратегия, а не просто проверява стойност.
    if debug_mode:
        # `debug_box` пази резултата от `slide.shapes.add_shape`, за да бъде проверен или използван в следващите стъпки вместо операцията да се повтори.
        debug_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(_px(left)), Inches(_px(top)), Inches(_px(width)), Inches(_px(height))
        )
        debug_box.fill.background()
        # `debug_box.line.color.rgb` пази резултата от `_rgb`, за да бъде проверен или използван в следващите стъпки вместо операцията да се повтори.
        debug_box.line.color.rgb = _rgb(theme.tokens.accent_primary)
        # `debug_box.line.width` пази резултата от `Pt`, за да бъде проверен или използван в следващите стъпки вместо операцията да се повтори.
        debug_box.line.width = Pt(1)
        _render_debug_label(slide, element, left, top, theme)

    # Обхождаме `element.children` като `child`, защото всеки елемент трябва да мине през една и съща pipeline стъпка.
    # Цикълът държи обработката еднаква за всеки елемент.
    for child in element.children:
        _render_layout_element(slide, child, theme, offset_x=left, offset_y=top, debug_mode=debug_mode)


def _render_layout_slide(prs: PptxPresentation, slide_data: LayoutedSlide, theme: ThemeDefinition) -> None:
    # Роля в pipeline-а: Това е вътрешна помощна стъпка: превежда semantic/layout информация към конкретни визуални обекти.
    # Входът идва през `prs` (PptxPresentation), `slide_data` (LayoutedSlide), `theme` (ThemeDefinition); имената показват каква част от контекста е собственост на тази стъпка.
    # Основните преходи навън са към `prs.slides.add_slide`, `page.background.fill.solid`, `_rgb`, `page.shapes.add_shape`; така се вижда кои отговорности функцията делегира.
    # Типовете в сигнатурата документират договора за caller-а и позволяват грешки да се хващат преди runtime.
    # Изходен договор: функцията не връща нов обект; ефектът ѝ е промяна на подадено състояние, файл или външна услуга.
    # `page` пази резултата от `prs.slides.add_slide`, за да бъде проверен или използван в следващите стъпки вместо операцията да се повтори.
    page = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT_INDEX])
    page.background.fill.solid()
    # `page.background.fill.fore_color.rgb` пази резултата от `_rgb`, за да бъде проверен или използван в следващите стъпки вместо операцията да се повтори.
    page.background.fill.fore_color.rgb = _rgb(theme.tokens.background)

    # `background` пази резултата от `page.shapes.add_shape`, за да бъде проверен или използван в следващите стъпки вместо операцията да се повтори.
    background = page.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT)
    background.fill.gradient()
    background.fill.gradient_angle = 135
    stops = background.fill.gradient_stops
    stops[0].position = 0.0
    # `stops[0].color.rgb` пази резултата от `_rgb`, за да бъде проверен или използван в следващите стъпки вместо операцията да се повтори.
    stops[0].color.rgb = _rgb(theme.tokens.background)
    stops[1].position = 1.0
    # `stops[1].color.rgb` пази резултата от `_rgb`, за да бъде проверен или използван в следващите стъпки вместо операцията да се повтори.
    stops[1].color.rgb = _rgb(theme.tokens.background_alt)
    background.line.fill.background()

    # `accent_position` пази резултата от `_component_style(theme, 'background').get`, за да бъде проверен или използван в следващите стъпки вместо операцията да се повтори.
    accent_position = _component_style(theme, "background").get("accent_position")
    # `accent` пази резултата от `page.shapes.add_shape`, за да бъде проверен или използван в следващите стъпки вместо операцията да се повтори.
    accent = page.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        SLIDE_WIDTH if accent_position == "top" else Inches(0.12),
        Inches(0.12) if accent_position == "top" else SLIDE_HEIGHT,
    )
    accent.fill.solid()
    # `accent.fill.fore_color.rgb` пази резултата от `_rgb`, за да бъде проверен или използван в следващите стъпки вместо операцията да се повтори.
    accent.fill.fore_color.rgb = _rgb(theme.tokens.accent_primary)
    accent.line.fill.background()

    # Обхождаме `slide_data.elements` като `element`, защото всеки елемент трябва да мине през една и съща pipeline стъпка.
    # Цикълът държи обработката еднаква за всеки елемент.
    for element in slide_data.elements:
        _render_layout_element(page, element, theme, debug_mode=slide_data.debug_mode)


def build_pptx_presentation(
    layouted_presentation: LayoutedPresentationDocument, theme: ThemeDefinition
) -> PptxPresentation:
    # Роля в pipeline-а: сглобява по-ниско ниво данни в обект, който следващият pipeline етап разбира директно.
    # Входът идва през `layouted_presentation` (LayoutedPresentationDocument), `theme` (ThemeDefinition); имената показват каква част от контекста е собственост на тази стъпка.
    # Основните преходи навън са към `PptxPresentation`, `_render_layout_slide`; така се вижда кои отговорности функцията делегира.
    # Типовете в сигнатурата документират договора за caller-а и позволяват грешки да се хващат преди runtime.
    # Изходен договор: `PptxPresentation`. Резултатът е част от последния rendering/export етап и вече е близо до крайния PPTX/PDF файл.
    # `pptx` пази резултата от `PptxPresentation`, за да бъде проверен или използван в следващите стъпки вместо операцията да се повтори.
    pptx = PptxPresentation()
    pptx.slide_width = SLIDE_WIDTH
    pptx.slide_height = SLIDE_HEIGHT

    # Обхождаме `layouted_presentation.slides` като `slide`, защото всеки елемент трябва да мине през една и съща pipeline стъпка.
    # Цикълът държи обработката еднаква за всеки елемент.
    for slide in layouted_presentation.slides:
        _render_layout_slide(pptx, slide, theme)

    return pptx
