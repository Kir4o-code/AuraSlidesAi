# Роля на модула: Legacy browser-based export път, използван като отделна стратегия.
# Чети коментарите като обяснение на причината за кода и връзката му със следващия слой, а не като буквален превод на Python синтаксиса.
import json
import logging
import os
import subprocess
from pathlib import Path

from PIL import Image
from pptx import Presentation as PptxPresentation
from pptx.util import Inches

from app.schemas.presentation import Presentation
from app.services.exporters.base import BaseExporter

logger = logging.getLogger(__name__)


class ScreenshotExporter(BaseExporter):
    # Роля на класа: Класът групира общо състояние и операции, които принадлежат на една pipeline отговорност.
    # Методите получават `self`, затова могат да споделят конфигурация и кеширани ресурси без глобални променливи.
    """
    Legacy fallback exporter that captures slides by screenshotting the React frontend.
    This is intentionally DOM-dependent and must remain opt-in only.
    """

    def __init__(self) -> None:
        # Роля в pipeline-а: обработва стъпката `init` като отделна отговорност, така че caller-ът да използва резултата без да познава вътрешните проверки и междинни стойности.
        # Входът идва през `self` (неуточнен тип); имената показват каква част от контекста е собственост на тази стъпка.
        # Основните преходи навън са към `os.getenv('FRONTEND_ORIGIN', 'http://localhost:3000').rstrip`, `Inches`, `os.getenv`, `Path`; така се вижда кои отговорности функцията делегира.
        # Типовете в сигнатурата документират договора за caller-а и позволяват грешки да се хващат преди runtime.
        # Изходен договор: функцията не връща нов обект; ефектът ѝ е промяна на подадено състояние, файл или външна услуга.
        # `self.backend_dir` пази резултата от `Path(__file__).resolve`, за да бъде проверен или използван в следващите стъпки вместо операцията да се повтори.
        self.backend_dir = Path(__file__).resolve().parents[3]
        self.output_dir = self.backend_dir / "generated"
        self.screenshot_dir = self.output_dir / "react_screenshots"
        self.export_data_dir = self.output_dir / "export_data"
        # `self.frontend_origin` пази резултата от `os.getenv('FRONTEND_ORIGIN', 'http://localhost:3000').rstrip`, за да бъде проверен или използван в следващите стъпки вместо операцията да се повтори.
        self.frontend_origin = os.getenv("FRONTEND_ORIGIN", "http://localhost:3000").rstrip("/")
        # `self.timeout` пази резултата от `os.getenv`, за да бъде проверен или използван в следващите стъпки вместо операцията да се повтори.
        self.timeout = int(os.getenv("REACT_EXPORT_TIMEOUT_SECONDS", "90"))
        # `self.enabled` пази резултата от `os.getenv('ENABLE_LEGACY_SCREENSHOT_EXPORT', 'false').strip().lower`, за да бъде проверен или използван в следващите стъпки вместо операцията да се повтори.
        self.enabled = os.getenv("ENABLE_LEGACY_SCREENSHOT_EXPORT", "false").strip().lower() in {
            "1",
            "true",
            "yes",
            "on",
        }

        self.slide_width_px = 1280
        self.slide_height_px = 720
        # `self.pptx_width` пази резултата от `Inches`, за да бъде проверен или използван в следващите стъпки вместо операцията да се повтори.
        self.pptx_width = Inches(13.333333)
        # `self.pptx_height` пази резултата от `Inches`, за да бъде проверен или използван в следващите стъпки вместо операцията да се повтори.
        self.pptx_height = Inches(7.5)

        # `self.chrome_candidates` пази резултата от `Path`, за да бъде проверен или използван в следващите стъпки вместо операцията да се повтори.
        self.chrome_candidates = (
            Path(os.getenv("CHROME_PATH", "")) if os.getenv("CHROME_PATH") else None,
            Path(r"C:\Program Files\Google\Chrome\Application\chrome.exe"),
            Path(r"C:\Program Files (x86)\Google\Chrome\Application\chrome.exe"),
            Path(r"C:\Program Files\Microsoft\Edge\Application\msedge.exe"),
            Path(r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe"),
        )

    def _find_browser(self) -> Path:
        # Роля в pipeline-а: Това е вътрешна помощна стъпка: обработва стъпката `find_browser` като отделна отговорност, така че caller-ът да използва резултата без да познава вътрешните проверки и междинни стойности.
        # Входът идва през `self` (неуточнен тип); имената показват каква част от контекста е собственост на тази стъпка.
        # Основните преходи навън са към `RuntimeError`, `candidate.exists`; така се вижда кои отговорности функцията делегира.
        # Типовете в сигнатурата документират договора за caller-а и позволяват грешки да се хващат преди runtime.
        # Изходен договор: `Path`. Резултатът е част от последния rendering/export етап и вече е близо до крайния PPTX/PDF файл.
        # Това условие е decision point: `not self.enabled`.
        # При вярно условие се активира `RuntimeError`; така този branch избира конкретна стратегия, а не просто проверява стойност.
        if not self.enabled:
            raise RuntimeError(
                "Legacy screenshot export is disabled. Set ENABLE_LEGACY_SCREENSHOT_EXPORT=true to use it."
            )
        # Обхождаме `self.chrome_candidates` като `candidate`, защото всеки елемент трябва да мине през една и съща pipeline стъпка.
        # Цикълът държи обработката еднаква за всеки елемент.
        for candidate in self.chrome_candidates:
            # Това условие е decision point: `candidate and candidate.exists()`.
            # Това е guard clause: при вярно условие вече имаме достатъчно надежден резултат (`candidate`) и прескачаме ненужната останала работа.
            if candidate and candidate.exists():
                return candidate
        raise RuntimeError("Chromium-based browser not found for screenshot export.")

    def _write_data(self, presentation: Presentation, asset_id: str) -> Path:
        # Роля в pipeline-а: Това е вътрешна помощна стъпка: обработва стъпката `write_data` като отделна отговорност, така че caller-ът да използва резултата без да познава вътрешните проверки и междинни стойности.
        # Входът идва през `self` (неуточнен тип), `presentation` (Presentation), `asset_id` (str); имената показват каква част от контекста е собственост на тази стъпка.
        # Основните преходи навън са към `self.export_data_dir.mkdir`, `prepare_export_bundle`, `path.write_text`, `json.dumps`; така се вижда кои отговорности функцията делегира.
        # Типовете в сигнатурата документират договора за caller-а и позволяват грешки да се хващат преди runtime.
        # Изходен договор: `Path`. Резултатът е част от последния rendering/export етап и вече е близо до крайния PPTX/PDF файл.
        from app.services.slide_generator import prepare_export_bundle

        self.export_data_dir.mkdir(parents=True, exist_ok=True)
        path = self.export_data_dir / f"{asset_id}.json"
        layouted_presentation, _ = prepare_export_bundle(presentation)
        path.write_text(
            json.dumps(
                {
                    "presentation": presentation.model_dump(mode="json"),
                    "layouted_presentation": layouted_presentation.model_dump(mode="json"),
                },
                ensure_ascii=False,
            ),
            encoding="utf-8",
        )
        return path

    def export_pptx(self, presentation: Presentation, theme: object, asset_id: str) -> str:
        # Роля в pipeline-а: превръща готовия layout в краен файл за клиента.
        # Входът идва през `self` (неуточнен тип), `presentation` (Presentation), `theme` (object), `asset_id` (str); имената показват каква част от контекста е собственост на тази стъпка.
        # Основните преходи навън са към `self.output_dir.mkdir`, `self.screenshot_dir.mkdir`, `self._write_data`, `self._capture_all_slides`; така се вижда кои отговорности функцията делегира.
        # Типовете в сигнатурата документират договора за caller-а и позволяват грешки да се хващат преди runtime.
        # Изходен договор: `str`. Резултатът е част от последния rendering/export етап и вече е близо до крайния PPTX/PDF файл.
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.screenshot_dir.mkdir(parents=True, exist_ok=True)
        self._write_data(presentation, asset_id)

        # `screenshot_paths` пази резултата от `self._capture_all_slides`, за да бъде проверен или използван в следващите стъпки вместо операцията да се повтори.
        screenshot_paths = self._capture_all_slides(presentation, asset_id)

        filename = f"{asset_id}.pptx"
        # `output_path` е крайното място във файловата система, което следващият слой може безопасно да използва.
        output_path = self.output_dir / filename

        # `deck` пази резултата от `PptxPresentation`, за да бъде проверен или използван в следващите стъпки вместо операцията да се повтори.
        deck = PptxPresentation()
        deck.slide_width = self.pptx_width
        deck.slide_height = self.pptx_height
        blank_layout = deck.slide_layouts[6]

        # Обхождаме `screenshot_paths` като `screenshot_path`, защото всеки елемент трябва да мине през една и съща pipeline стъпка.
        # Цикълът държи обработката еднаква за всеки елемент.
        for screenshot_path in screenshot_paths:
            # `slide` пази резултата от `deck.slides.add_slide`, за да бъде проверен или използван в следващите стъпки вместо операцията да се повтори.
            slide = deck.slides.add_slide(blank_layout)
            slide.shapes.add_picture(str(screenshot_path), 0, 0, width=self.pptx_width, height=self.pptx_height)

        deck.save(output_path)
        return filename

    def export_pdf(self, presentation: Presentation, theme: object, asset_id: str) -> str:
        # Роля в pipeline-а: превръща готовия layout в краен файл за клиента.
        # Входът идва през `self` (неуточнен тип), `presentation` (Presentation), `theme` (object), `asset_id` (str); имената показват каква част от контекста е собственост на тази стъпка.
        # Основните преходи навън са към `self.output_dir.mkdir`, `self.screenshot_dir.mkdir`, `self._write_data`, `self._capture_all_slides`; така се вижда кои отговорности функцията делегира.
        # Типовете в сигнатурата документират договора за caller-а и позволяват грешки да се хващат преди runtime.
        # Изходен договор: `str`. Резултатът е част от последния rendering/export етап и вече е близо до крайния PPTX/PDF файл.
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.screenshot_dir.mkdir(parents=True, exist_ok=True)
        self._write_data(presentation, asset_id)

        # `screenshot_paths` пази резултата от `self._capture_all_slides`, за да бъде проверен или използван в следващите стъпки вместо операцията да се повтори.
        screenshot_paths = self._capture_all_slides(presentation, asset_id)

        filename = f"{asset_id}.pdf"
        # `output_path` е крайното място във файловата система, което следващият слой може безопасно да използва.
        output_path = self.output_dir / filename

        # `images` пази резултата от `Image.open(p).convert`, за да бъде проверен или използван в следващите стъпки вместо операцията да се повтори.
        # Comprehension синтаксисът комбинира обхождане и филтриране в една стойност; резултатът съдържа само елементите, минали условието.
        images = [Image.open(p).convert("RGB") for p in screenshot_paths]
        # Това условие е decision point: `not images`.
        # При вярно условие се активира `RuntimeError`; така този branch избира конкретна стратегия, а не просто проверява стойност.
        if not images:
            raise RuntimeError("No screenshots captured for PDF.")

        # Тук започва контролирана рискова зона: външна услуга, parsing, filesystem или rendering може да се провали.
        # `try/except` превръща техническите грешки () в предвидимо поведение за горния слой.
        try:
            images[0].save(output_path, "PDF", resolution=144.0, save_all=True, append_images=images[1:])
        finally:
            # Обхождаме `images` като `image`, защото всеки елемент трябва да мине през една и съща pipeline стъпка.
            # Цикълът държи обработката еднаква за всеки елемент.
            for image in images:
                image.close()

        return filename

    def _capture_all_slides(self, presentation: Presentation, asset_id: str) -> list[Path]:
        # Роля в pipeline-а: Това е вътрешна помощна стъпка: обработва стъпката `capture_all_slides` като отделна отговорност, така че caller-ът да използва резултата без да познава вътрешните проверки и междинни стойности.
        # Входът идва през `self` (неуточнен тип), `presentation` (Presentation), `asset_id` (str); имената показват каква част от контекста е собственост на тази стъпка.
        # Основните преходи навън са към `self._find_browser`, `subprocess.run`; така се вижда кои отговорности функцията делегира.
        # Типовете в сигнатурата документират договора за caller-а и позволяват грешки да се хващат преди runtime.
        # Изходен договор: `list[Path]`. Резултатът е част от последния rendering/export етап и вече е близо до крайния PPTX/PDF файл.
        # `browser` пази резултата от `self._find_browser`, за да бъде проверен или използван в следващите стъпки вместо операцията да се повтори.
        browser = self._find_browser()
        paths: list[Path] = []
        # Обхождаме `range(len(presentation.slides))` като `index`, защото всеки елемент трябва да мине през една и съща pipeline стъпка.
        # Цикълът държи обработката еднаква за всеки елемент.
        for index in range(len(presentation.slides)):
            # `output_path` е крайното място във файловата система, което следващият слой може безопасно да използва.
            output_path = self.screenshot_dir / f"{asset_id}-slide-{index + 1:02d}.png"
            url = f"{self.frontend_origin}/export/{asset_id}?slide={index}"

            command = [
                str(browser),
                "--headless=new",
                "--disable-gpu",
                "--hide-scrollbars",
                f"--window-size={self.slide_width_px},{self.slide_height_px}",
                f"--screenshot={output_path}",
                "--virtual-time-budget=3000",
                url,
            ]

            subprocess.run(command, timeout=self.timeout, check=True, capture_output=True)
            paths.append(output_path)
        return paths
