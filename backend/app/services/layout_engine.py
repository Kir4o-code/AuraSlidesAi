from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation as PptxPresentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from app.semantic.contracts import Alignment, LayoutElement, LayoutElementKind, LayoutedPresentationDocument, ThemeDefinition


SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
BLANK_LAYOUT_INDEX = 6


def _rgb(value: str) -> RGBColor:
    cleaned = value.strip().lstrip("#")
    if len(cleaned) == 3:
        cleaned = "".join(char * 2 for char in cleaned)
    if len(cleaned) != 6:
        return RGBColor(37, 99, 235)
    return RGBColor(*(int(cleaned[index : index + 2], 16) for index in (0, 2, 4)))


PX_TO_INCH = 1.0 / 96.0


def _px(value: int) -> float:
    return value * PX_TO_INCH


def _component_style(theme: ThemeDefinition, component: str) -> dict:
    return theme.tokens.component_styles.get(component, {})


def _add_fitted_picture(slide, image_path: Path, left: int, top: int, width: int, height: int, fit: str) -> None:
    with Image.open(image_path) as image:
        source_width, source_height = image.size
    source_ratio = source_width / max(source_height, 1)
    frame_ratio = width / max(height, 1)
    if fit == "contain":
        if source_ratio > frame_ratio:
            render_width = width
            render_height = width / source_ratio
        else:
            render_height = height
            render_width = height * source_ratio
        slide.shapes.add_picture(
            str(image_path),
            Inches(_px(left + ((width - render_width) / 2))),
            Inches(_px(top + ((height - render_height) / 2))),
            Inches(_px(render_width)),
            Inches(_px(render_height)),
        )
        return
    picture = slide.shapes.add_picture(str(image_path), Inches(_px(left)), Inches(_px(top)), Inches(_px(width)), Inches(_px(height)))
    if source_ratio > frame_ratio:
        crop = max(0.0, (1 - (frame_ratio / source_ratio)) / 2)
        picture.crop_left = crop
        picture.crop_right = crop
    elif source_ratio < frame_ratio:
        crop = max(0.0, (1 - (source_ratio / frame_ratio)) / 2)
        picture.crop_top = crop
        picture.crop_bottom = crop


def _add_semantic_icon(slide, icon: str, left: float, top: float, size: float, theme: ThemeDefinition) -> None:
    accent = _rgb(theme.tokens.accent_primary)
    soft = _rgb(theme.tokens.accent_secondary)

    def add_shape(kind, x: float, y: float, width: float, height: float, *, filled: bool = False):
        shape = slide.shapes.add_shape(kind, Inches(_px(x)), Inches(_px(y)), Inches(_px(width)), Inches(_px(height)))
        if filled:
            shape.fill.solid()
            shape.fill.fore_color.rgb = accent
        else:
            shape.fill.background()
        shape.line.color.rgb = accent
        shape.line.width = Pt(1)
        return shape

    def add_line(x1: float, y1: float, x2: float, y2: float) -> None:
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(_px(x1)), Inches(_px(y1)), Inches(_px(x2)), Inches(_px(y2)))
        line.line.color.rgb = accent
        line.line.width = Pt(1)

    container = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(_px(left)), Inches(_px(top)), Inches(_px(size)), Inches(_px(size)))
    container.fill.solid()
    container.fill.fore_color.rgb = soft
    container.line.color.rgb = accent
    container.line.width = Pt(0.8)
    x = left + (size * 0.23)
    y = top + (size * 0.23)
    inner = size * 0.54
    if icon == "chart":
        for index, factor in enumerate((0.42, 0.72, 1.0)):
            bar_height = inner * factor
            add_shape(MSO_SHAPE.RECTANGLE, x + (index * inner * 0.34), y + inner - bar_height, inner * 0.16, bar_height, filled=True)
    elif icon == "bolt":
        add_shape(MSO_SHAPE.LIGHTNING_BOLT, x, y, inner, inner, filled=True)
    elif icon == "star":
        add_shape(MSO_SHAPE.STAR_5_POINT, x, y, inner, inner, filled=True)
    elif icon == "heart":
        add_shape(MSO_SHAPE.HEART, x, y, inner, inner, filled=True)
    elif icon in {"person", "users"}:
        add_shape(MSO_SHAPE.OVAL, x + inner * 0.32, y, inner * 0.36, inner * 0.36, filled=True)
        add_shape(MSO_SHAPE.OVAL, x + inner * 0.14, y + inner * 0.42, inner * 0.72, inner * 0.5)
    elif icon == "clock":
        add_shape(MSO_SHAPE.OVAL, x, y, inner, inner)
        add_line(x + inner / 2, y + inner / 2, x + inner / 2, y + inner * 0.2)
        add_line(x + inner / 2, y + inner / 2, x + inner * 0.74, y + inner * 0.58)
    elif icon == "eye":
        add_shape(MSO_SHAPE.OVAL, x, y + inner * 0.16, inner, inner * 0.68)
        add_shape(MSO_SHAPE.OVAL, x + inner * 0.34, y + inner * 0.34, inner * 0.32, inner * 0.32, filled=True)
    elif icon == "book":
        add_shape(MSO_SHAPE.RECTANGLE, x, y + inner * 0.08, inner * 0.46, inner * 0.84)
        add_shape(MSO_SHAPE.RECTANGLE, x + inner * 0.54, y + inner * 0.08, inner * 0.46, inner * 0.84)
    elif icon in {"map", "route"}:
        add_line(x + inner * 0.12, y + inner * 0.82, x + inner * 0.48, y + inner * 0.28)
        add_line(x + inner * 0.48, y + inner * 0.28, x + inner * 0.88, y + inner * 0.68)
        for px, py in ((0.12, 0.82), (0.48, 0.28), (0.88, 0.68)):
            add_shape(MSO_SHAPE.OVAL, x + inner * px - 2, y + inner * py - 2, 4, 4, filled=True)
    elif icon == "search":
        add_shape(MSO_SHAPE.OVAL, x, y, inner * 0.68, inner * 0.68)
        add_line(x + inner * 0.56, y + inner * 0.56, x + inner, y + inner)
    elif icon == "film":
        add_shape(MSO_SHAPE.RECTANGLE, x, y + inner * 0.12, inner, inner * 0.76)
        add_line(x + inner * 0.25, y + inner * 0.12, x + inner * 0.25, y + inner * 0.88)
        add_line(x + inner * 0.75, y + inner * 0.12, x + inner * 0.75, y + inner * 0.88)
    elif icon == "shield":
        add_shape(MSO_SHAPE.PENTAGON, x, y, inner, inner, filled=True)
    elif icon == "home":
        add_shape(MSO_SHAPE.PENTAGON, x, y, inner, inner)
        add_shape(MSO_SHAPE.RECTANGLE, x + inner * 0.38, y + inner * 0.55, inner * 0.24, inner * 0.45)
    elif icon == "idea":
        add_shape(MSO_SHAPE.SUN, x, y, inner, inner)
    else:
        add_shape(MSO_SHAPE.OVAL, x, y, inner, inner)
        add_shape(MSO_SHAPE.OVAL, x + inner * 0.24, y + inner * 0.24, inner * 0.52, inner * 0.52)


def _semantic_font_name(theme: ThemeDefinition, region: str, kind: LayoutElementKind) -> str:
    if kind in {LayoutElementKind.QUOTE, LayoutElementKind.STATISTIC}:
        return theme.tokens.fonts.heading or theme.tokens.fonts.body
    if any(token in region for token in ("title", "quote", "attribution", "heading")):
        return theme.tokens.fonts.heading or theme.tokens.fonts.body
    return theme.tokens.fonts.body or theme.tokens.fonts.heading


def _render_debug_label(slide, element: LayoutElement, left: float, top: float, theme: ThemeDefinition) -> None:
    box = slide.shapes.add_textbox(Inches(_px(left + 4)), Inches(_px(top + 4)), Inches(max(_px(min(element.width, 200)), 0.4)), Inches(0.22))
    frame = box.text_frame
    frame.clear()
    run = frame.paragraphs[0].add_run()
    run.text = f"{element.region} · {element.x},{element.y} {element.width}x{element.height}"
    run.font.name = theme.tokens.fonts.mono or theme.tokens.fonts.body
    run.font.size = Pt(7)
    run.font.color.rgb = _rgb(theme.tokens.accent_primary)


def _render_layout_element(slide, element: LayoutElement, theme: ThemeDefinition, *, offset_x: int = 0, offset_y: int = 0, debug_mode: bool = False) -> None:
    left = offset_x + element.x
    top = offset_y + element.y
    width = element.width
    height = element.height

    if element.kind == LayoutElementKind.PANEL:
        panel_shape = MSO_SHAPE.RECTANGLE if _component_style(theme, "panel").get("style") == "square" else MSO_SHAPE.ROUNDED_RECTANGLE
        panel = slide.shapes.add_shape(panel_shape, Inches(_px(left)), Inches(_px(top)), Inches(_px(width)), Inches(_px(height)))
        panel.fill.solid()
        panel.fill.fore_color.rgb = _rgb(theme.tokens.surface)
        panel.line.color.rgb = _rgb(theme.tokens.border)
        panel.line.width = Pt(1)
    elif element.kind == LayoutElementKind.IMAGE:
        local_path = element.content.get("local_path") if isinstance(element.content, dict) else None
        image_path = Path(local_path) if local_path else None
        image_style = _component_style(theme, "image")
        inset = max(0, int(image_style.get("frame_inset") or 0))
        frame_shape = MSO_SHAPE.RECTANGLE if _component_style(theme, "panel").get("style") == "square" else MSO_SHAPE.ROUNDED_RECTANGLE
        frame_shape = slide.shapes.add_shape(frame_shape, Inches(_px(left)), Inches(_px(top)), Inches(_px(width)), Inches(_px(height)))
        frame_shape.fill.solid()
        frame_shape.fill.fore_color.rgb = _rgb(theme.tokens.background_alt)
        frame_shape.line.color.rgb = _rgb(theme.tokens.border)
        frame_shape.line.width = Pt(1)
        if image_path and image_path.exists():
            _add_fitted_picture(
                slide,
                image_path,
                left + inset,
                top + inset,
                max(1, width - (inset * 2)),
                max(1, height - (inset * 2)),
                str(element.content.get("fit") or image_style.get("fit") or "cover"),
            )
        else:
            placeholder = slide.shapes.add_textbox(Inches(_px(left + inset)), Inches(_px(top + inset)), Inches(_px(max(1, width - (inset * 2)))), Inches(_px(max(1, height - (inset * 2)))))
            frame = placeholder.text_frame
            frame.word_wrap = True
            frame.clear()
            run = frame.paragraphs[0].add_run()
            run.text = element.content.get("prompt") or element.text or "Image"
            run.font.name = theme.tokens.fonts.body
            run.font.size = Pt(max(10, (element.font_size or 18) * 0.7))
            run.font.color.rgb = _rgb(theme.tokens.text_primary)
    elif element.kind != LayoutElementKind.BULLET_LIST:
        text_left = left
        text_width = width
        text_top = top
        text_height = height
        if element.kind == LayoutElementKind.BULLET_ITEM:
            if _component_style(theme, "bullet").get("style") == "lines":
                line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(_px(left)), Inches(_px(top)), Inches(_px(3)), Inches(_px(height)))
                line.fill.solid()
                line.fill.fore_color.rgb = _rgb(theme.tokens.accent_primary)
                line.line.fill.background()
            else:
                card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(_px(left)), Inches(_px(top)), Inches(_px(width)), Inches(_px(height)))
                card.fill.solid()
                card.fill.fore_color.rgb = _rgb(theme.tokens.surface)
                card.line.color.rgb = _rgb(theme.tokens.border)
                card.line.width = Pt(0.8)
            marker_size = min(28, max(20, height - 12))
            _add_semantic_icon(slide, str(element.content.get("icon") or "target"), left + 10, top + ((height - marker_size) / 2), marker_size, theme)
            text_left += 50
            text_width -= 58
            text_top += 5
            text_height -= 10
        textbox = slide.shapes.add_textbox(Inches(_px(text_left)), Inches(_px(text_top)), Inches(_px(max(1, text_width))), Inches(_px(max(1, text_height))))
        frame = textbox.text_frame
        frame.word_wrap = element.wrap
        frame.auto_size = MSO_AUTO_SIZE.NONE
        frame.vertical_anchor = MSO_ANCHOR.TOP
        frame.clear()
        frame.margin_left = Inches(0.02)
        frame.margin_right = Inches(0.02)
        frame.margin_top = Inches(0.02)
        frame.margin_bottom = Inches(0.02)
        paragraph = frame.paragraphs[0]
        paragraph.alignment = {
            Alignment.START: PP_ALIGN.LEFT,
            Alignment.CENTER: PP_ALIGN.CENTER,
            Alignment.END: PP_ALIGN.RIGHT,
        }[element.align]
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(0)
        paragraph.line_spacing = min(element.line_height or 1.18, 1.18)
        run = paragraph.add_run()
        run.text = element.text or ""
        run.font.name = _semantic_font_name(theme, element.region, element.kind)
        if element.font_size:
            run.font.size = Pt(round(element.font_size * 0.75, 1))
        run.font.color.rgb = _rgb(theme.tokens.text_secondary if any(token in element.region for token in ("subtitle", "notes", "attribution")) else theme.tokens.text_primary)

    if debug_mode:
        debug_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(_px(left)), Inches(_px(top)), Inches(_px(width)), Inches(_px(height)))
        debug_box.fill.background()
        debug_box.line.color.rgb = _rgb(theme.tokens.accent_primary)
        debug_box.line.width = Pt(1)
        _render_debug_label(slide, element, left, top, theme)

    for child in element.children:
        _render_layout_element(slide, child, theme, offset_x=left, offset_y=top, debug_mode=debug_mode)


def _render_layout_slide(prs: PptxPresentation, slide_data: LayoutedSlide, theme: ThemeDefinition) -> None:
    page = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT_INDEX])
    page.background.fill.solid()
    page.background.fill.fore_color.rgb = _rgb(theme.tokens.background)

    background = page.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT)
    background.fill.gradient()
    background.fill.gradient_angle = 135
    stops = background.fill.gradient_stops
    stops[0].position = 0.0
    stops[0].color.rgb = _rgb(theme.tokens.background)
    stops[1].position = 1.0
    stops[1].color.rgb = _rgb(theme.tokens.background_alt)
    background.line.fill.background()

    accent_position = _component_style(theme, "background").get("accent_position")
    accent = page.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        SLIDE_WIDTH if accent_position == "top" else Inches(0.12),
        Inches(0.12) if accent_position == "top" else SLIDE_HEIGHT,
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = _rgb(theme.tokens.accent_primary)
    accent.line.fill.background()

    for element in slide_data.elements:
        _render_layout_element(page, element, theme, debug_mode=slide_data.debug_mode)


def build_pptx_presentation(layouted_presentation: LayoutedPresentationDocument, theme: ThemeDefinition) -> PptxPresentation:
    pptx = PptxPresentation()
    pptx.slide_width = SLIDE_WIDTH
    pptx.slide_height = SLIDE_HEIGHT

    for slide in layouted_presentation.slides:
        _render_layout_slide(pptx, slide, theme)

    return pptx
