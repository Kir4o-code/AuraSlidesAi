import json
import logging
import os
import subprocess
from pathlib import Path
from time import perf_counter

from PIL import Image
from pptx import Presentation as PptxPresentation
from pptx.util import Inches

from app.schemas.presentation import Presentation


class ReactExportError(Exception):
    pass


logger = logging.getLogger(__name__)
BACKEND_DIR = Path(__file__).resolve().parents[2]
OUTPUT_DIR = BACKEND_DIR / "generated"
EXPORT_DATA_DIR = OUTPUT_DIR / "export_data"
SCREENSHOT_DIR = OUTPUT_DIR / "react_screenshots"
FRONTEND_ORIGIN = os.getenv("FRONTEND_ORIGIN", "http://localhost:3000").rstrip("/")
REACT_EXPORT_TIMEOUT_SECONDS = int(os.getenv("REACT_EXPORT_TIMEOUT_SECONDS", "90"))
SLIDE_WIDTH_PX = 1280
SLIDE_HEIGHT_PX = 720
PPTX_WIDTH = Inches(13.333333)
PPTX_HEIGHT = Inches(7.5)
CHROME_CANDIDATES = (
    Path(os.getenv("CHROME_PATH", "")) if os.getenv("CHROME_PATH") else None,
    Path(r"C:\Program Files\Google\Chrome\Application\chrome.exe"),
    Path(r"C:\Program Files (x86)\Google\Chrome\Application\chrome.exe"),
    Path(r"C:\Program Files\Microsoft\Edge\Application\msedge.exe"),
    Path(r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe"),
)


def _find_browser_executable() -> Path:
    for candidate in CHROME_CANDIDATES:
        if candidate and candidate.exists():
            return candidate
    raise ReactExportError("Chrome or Edge was not found. Set CHROME_PATH or install a Chromium-based browser.")


def _run_browser(command: list[str], *, action: str) -> None:
    try:
        result = subprocess.run(
            command,
            cwd=BACKEND_DIR,
            capture_output=True,
            text=True,
            timeout=REACT_EXPORT_TIMEOUT_SECONDS,
            check=False,
        )
    except subprocess.TimeoutExpired as exc:
        raise ReactExportError(
            f"{action} timed out after {REACT_EXPORT_TIMEOUT_SECONDS}s. "
            f"Make sure the frontend is running at {FRONTEND_ORIGIN}."
        ) from exc
    if result.stdout.strip():
        logger.info("%s stdout:\n%s", action, result.stdout.strip())
    if result.stderr.strip():
        logger.warning("%s stderr:\n%s", action, result.stderr.strip())
    if result.returncode != 0:
        raise ReactExportError(f"{action} failed with exit code {result.returncode}.")


def _write_export_data(presentation: Presentation, asset_id: str) -> Path:
    EXPORT_DATA_DIR.mkdir(parents=True, exist_ok=True)
    output_path = EXPORT_DATA_DIR / f"{asset_id}.json"
    output_path.write_text(
        json.dumps(presentation.model_dump(mode="json"), ensure_ascii=False),
        encoding="utf-8",
    )
    return output_path


def _base_chrome_command(browser: Path) -> list[str]:
    return [
        str(browser),
        "--headless=new",
        "--disable-gpu",
        "--disable-background-networking",
        "--disable-extensions",
        "--no-first-run",
        "--no-default-browser-check",
        "--hide-scrollbars",
        "--allow-file-access-from-files",
        "--run-all-compositor-stages-before-draw",
        "--virtual-time-budget=3000",
    ]


def _capture_slide(browser: Path, asset_id: str, slide_index: int, output_path: Path) -> None:
    url = f"{FRONTEND_ORIGIN}/export/{asset_id}?slide={slide_index}"
    command = [
        *_base_chrome_command(browser),
        f"--window-size={SLIDE_WIDTH_PX},{SLIDE_HEIGHT_PX}",
        f"--screenshot={output_path}",
        url,
    ]
    logger.info("React slide screenshot starting. url=%s output=%s", url, output_path)
    _run_browser(command, action=f"React slide {slide_index + 1} screenshot")
    if not output_path.exists() or output_path.stat().st_size == 0:
        raise ReactExportError(f"React slide screenshot {slide_index + 1} was not created.")


def _build_pptx_from_screenshots(screenshot_paths: list[Path], output_path: Path, title: str) -> None:
    deck = PptxPresentation()
    deck.slide_width = PPTX_WIDTH
    deck.slide_height = PPTX_HEIGHT
    deck.core_properties.title = title
    deck.core_properties.subject = "AuraSlidesAi React-rendered presentation export"
    deck.core_properties.author = "AuraSlidesAi"
    blank_layout = deck.slide_layouts[6]

    for screenshot_path in screenshot_paths:
        slide = deck.slides.add_slide(blank_layout)
        slide.shapes.add_picture(str(screenshot_path), 0, 0, width=PPTX_WIDTH, height=PPTX_HEIGHT)

    deck.save(output_path)
    if not output_path.exists() or output_path.stat().st_size == 0:
        raise ReactExportError("React PPTX export finished without creating a PPTX file.")


def _build_pdf_from_screenshots(screenshot_paths: list[Path], output_path: Path) -> None:
    images: list[Image.Image] = []
    try:
        for screenshot_path in screenshot_paths:
            images.append(Image.open(screenshot_path).convert("RGB"))
        if not images:
            raise ReactExportError("No slide screenshots were available for PDF export.")
        first, rest = images[0], images[1:]
        first.save(output_path, "PDF", resolution=144.0, save_all=True, append_images=rest)
    finally:
        for image in images:
            image.close()

    if not output_path.exists() or output_path.stat().st_size == 0:
        raise ReactExportError("React PDF export finished without creating a PDF file.")


def build_react_presentation_exports(presentation: Presentation, asset_id: str) -> tuple[str, str]:
    started_at = perf_counter()
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    SCREENSHOT_DIR.mkdir(parents=True, exist_ok=True)
    export_data_path = _write_export_data(presentation, asset_id)
    browser = _find_browser_executable()

    pptx_name = f"{asset_id}.pptx"
    pdf_name = f"{asset_id}.pdf"
    pptx_path = OUTPUT_DIR / pptx_name
    pdf_path = OUTPUT_DIR / pdf_name
    screenshot_paths = [
        SCREENSHOT_DIR / f"{asset_id}-slide-{index + 1:02d}.png"
        for index in range(len(presentation.slides))
    ]

    logger.info(
        "React export starting. asset_id=%s slides=%s data=%s frontend=%s",
        asset_id,
        len(presentation.slides),
        export_data_path,
        FRONTEND_ORIGIN,
    )
    for index, screenshot_path in enumerate(screenshot_paths):
        _capture_slide(browser, asset_id, index, screenshot_path)
    _build_pptx_from_screenshots(screenshot_paths, pptx_path, presentation.title)
    _build_pdf_from_screenshots(screenshot_paths, pdf_path)

    logger.info(
        "React export complete in %.2fs. pptx=%s pdf=%s",
        perf_counter() - started_at,
        pptx_name,
        pdf_name,
    )
    return pptx_name, pdf_name
