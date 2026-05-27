import json
import logging
import os
import subprocess
from pathlib import Path

from PIL import Image
from pptx import Presentation as PptxPresentation
from pptx.util import Inches

from app.schemas.presentation import Presentation
from app.services.exporters.base import BaseExporter

logger = logging.getLogger(__name__)

class ScreenshotExporter(BaseExporter):
    """
    Legacy fallback exporter that captures slides by screenshotting the React frontend.
    This is intentionally DOM-dependent and must remain opt-in only.
    """
    
    def __init__(self):
        self.backend_dir = Path(__file__).resolve().parents[3]
        self.output_dir = self.backend_dir / "generated"
        self.screenshot_dir = self.output_dir / "react_screenshots"
        self.export_data_dir = self.output_dir / "export_data"
        self.frontend_origin = os.getenv("FRONTEND_ORIGIN", "http://localhost:3000").rstrip("/")
        self.timeout = int(os.getenv("REACT_EXPORT_TIMEOUT_SECONDS", "90"))
        self.enabled = os.getenv("ENABLE_LEGACY_SCREENSHOT_EXPORT", "false").strip().lower() in {"1", "true", "yes", "on"}
        
        self.slide_width_px = 1280
        self.slide_height_px = 720
        self.pptx_width = Inches(13.333333)
        self.pptx_height = Inches(7.5)
        
        self.chrome_candidates = (
            Path(os.getenv("CHROME_PATH", "")) if os.getenv("CHROME_PATH") else None,
            Path(r"C:\Program Files\Google\Chrome\Application\chrome.exe"),
            Path(r"C:\Program Files (x86)\Google\Chrome\Application\chrome.exe"),
            Path(r"C:\Program Files\Microsoft\Edge\Application\msedge.exe"),
            Path(r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe"),
        )

    def _find_browser(self) -> Path:
        if not self.enabled:
            raise RuntimeError(
                "Legacy screenshot export is disabled. Set ENABLE_LEGACY_SCREENSHOT_EXPORT=true to use it."
            )
        for candidate in self.chrome_candidates:
            if candidate and candidate.exists():
                return candidate
        raise RuntimeError("Chromium-based browser not found for screenshot export.")

    def _write_data(self, presentation: Presentation, asset_id: str) -> Path:
        self.export_data_dir.mkdir(parents=True, exist_ok=True)
        path = self.export_data_dir / f"{asset_id}.json"
        path.write_text(json.dumps(presentation.model_dump(mode="json"), ensure_ascii=False), encoding="utf-8")
        return path

    def export_pptx(self, presentation: Presentation, asset_id: str) -> str:
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.screenshot_dir.mkdir(parents=True, exist_ok=True)
        self._write_data(presentation, asset_id)
        
        screenshot_paths = self._capture_all_slides(presentation, asset_id)
        
        filename = f"{asset_id}.pptx"
        output_path = self.output_dir / filename
        
        deck = PptxPresentation()
        deck.slide_width = self.pptx_width
        deck.slide_height = self.pptx_height
        blank_layout = deck.slide_layouts[6]

        for screenshot_path in screenshot_paths:
            slide = deck.slides.add_slide(blank_layout)
            slide.shapes.add_picture(str(screenshot_path), 0, 0, width=self.pptx_width, height=self.pptx_height)

        deck.save(output_path)
        return filename

    def export_pdf(self, presentation: Presentation, asset_id: str) -> str:
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.screenshot_dir.mkdir(parents=True, exist_ok=True)
        self._write_data(presentation, asset_id)
        
        screenshot_paths = self._capture_all_slides(presentation, asset_id)
        
        filename = f"{asset_id}.pdf"
        output_path = self.output_dir / filename
        
        images = [Image.open(p).convert("RGB") for p in screenshot_paths]
        if not images:
            raise RuntimeError("No screenshots captured for PDF.")
            
        images[0].save(output_path, "PDF", resolution=144.0, save_all=True, append_images=images[1:])
        for img in images: img.close()
        
        return filename

    def _capture_all_slides(self, presentation: Presentation, asset_id: str) -> list[Path]:
        browser = self._find_browser()
        paths = []
        for index in range(len(presentation.slides)):
            output_path = self.screenshot_dir / f"{asset_id}-slide-{index + 1:02d}.png"
            url = f"{self.frontend_origin}/export/{asset_id}?slide={index}"
            
            command = [
                str(browser),
                "--headless=new",
                "--disable-gpu",
                "--hide-scrollbars",
                f"--window-size={self.slide_width_px},{self.slide_height_px}",
                f"--screenshot={output_path}",
                "--virtual-time-budget=3000",
                url
            ]
            
            subprocess.run(command, timeout=self.timeout, check=True, capture_output=True)
            paths.append(output_path)
        return paths
